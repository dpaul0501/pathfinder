"""Build PathFinder 4-slide pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x05, 0x96, 0x69)   # emerald
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT= RGBColor(0xED, 0xE9, 0xFE)
RED         = RGBColor(0xDC, 0x26, 0x26)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x0F, 0x17, 0x2A)
SLATE       = RGBColor(0x33, 0x41, 0x55)
MUTED       = RGBColor(0x64, 0x74, 0x8B)
BG          = RGBColor(0xF0, 0xF4, 0xF8)
SURFACE     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER      = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────

def bg_rect(slide, color=BG):
    sh = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def add_rect(slide, l, t, w, h, fill, line_color=None, radius=False):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Calibri", italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb

def add_label(slide, text, l, t, color=GREEN):
    """Small mono uppercase eyebrow label."""
    add_text(slide, text, l, t, 8, 0.3, size=9, bold=True, color=color,
             font="Courier New")

def accent_bar(slide, l, t, w=0.05, h=0.35, color=GREEN):
    add_rect(slide, l, t, w, h, fill=color)

def stat_tile(slide, l, t, w, h, val, label, val_color, bg_color, border_color, val_size=28):
    add_rect(slide, l, t, w, h, fill=bg_color, line_color=border_color)
    add_text(slide, val,   l+0.15, t+0.12, w-0.3, 0.55,
             size=val_size, bold=True, color=val_color, align=PP_ALIGN.LEFT)
    add_text(slide, label, l+0.15, t+0.62, w-0.3, 0.5,
             size=9, color=MUTED, align=PP_ALIGN.LEFT, wrap=True)

def pipe_node(slide, l, t, icon, name, tech, active=False):
    bg = GREEN_LIGHT if active else BORDER
    bc = GREEN if active else BORDER
    add_rect(slide, l, t, 1.45, 1.1, fill=bg, line_color=bc)
    add_text(slide, icon, l+0.55, t+0.08, 0.5, 0.4, size=18, align=PP_ALIGN.CENTER)
    add_text(slide, name, l+0.08, t+0.5,  1.29, 0.3, size=9.5, bold=True,
             color=DARK, align=PP_ALIGN.CENTER, font="Calibri")
    add_rect(slide, l+0.28, t+0.82, 0.9, 0.2, fill=GREEN_LIGHT if active else PURPLE_LIGHT)
    add_text(slide, tech, l+0.28, t+0.82, 0.9, 0.2, size=7.5,
             color=GREEN if active else PURPLE, align=PP_ALIGN.CENTER, font="Courier New")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
bg_rect(s1, BG)

# Left panel — dark story panel
add_rect(s1, 0, 0, 5.2, 7.5, fill=DARK)

# Eyebrow
add_label(s1, "HSI × WIT REGATTA 2026  ·  TRACK 2: YOUTH EMPOWERMENT", 0.4, 0.38, color=GREEN)

# Main headline
add_text(s1, "The door exists.", 0.4, 0.78, 4.4, 0.7,
         size=38, bold=True, color=WHITE, font="Calibri")
add_text(s1, "Marisol can't\nfind it.", 0.4, 1.42, 4.4, 1.4,
         size=38, bold=True, color=GREEN, font="Calibri")

# Story
story = (
    "She's 17. Latina. South Seattle.\n"
    "She wants to learn to code.\n\n"
    "She Googles 'free tech programs Seattle.'\n"
    "Gets 12 websites. 3 phone numbers.\n"
    "No clear next step.\n\n"
    "She gives up.\n\n"
    "The programs exist.\n"
    "She just can't find them."
)
add_text(s1, story, 0.4, 3.0, 4.4, 3.8,
         size=13, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

# Slide number
add_text(s1, "01 / 04", 0.4, 7.1, 1.5, 0.3, size=9,
         color=MUTED, font="Courier New")

# Right panel — stats
add_label(s1, "THE SCALE OF THE PROBLEM", 5.6, 0.5, color=MUTED)

# 4 stat tiles
stat_tile(s1, 5.5, 0.95, 3.6, 1.2,
          "~80,000", "WA youth ages 16–24 not in\neducation, employment, or training",
          GREEN, GREEN_LIGHT, GREEN)

stat_tile(s1, 5.5, 2.35, 3.6, 1.2,
          "30 points", "Latino credential gap vs. statewide\n(32% vs 62% post-secondary attainment)",
          PURPLE, PURPLE_LIGHT, PURPLE)

stat_tile(s1, 5.5, 3.75, 3.6, 1.2,
          "47 programs", "Free WA tech programs scattered\nacross 12 agencies, 5 nonprofits, 3 websites",
          AMBER, AMBER_LIGHT, AMBER)

stat_tile(s1, 5.5, 5.15, 3.6, 1.2,
          "0", "Unified navigation layers\nexisting before today",
          RED, RGBColor(0xFE, 0xE2, 0xE2), RED)

# WA Board of Ed quote
add_rect(s1, 5.5, 6.5, 7.6, 0.78, fill=BORDER)
add_text(s1,
         '"Policy has been developed in silos, without sufficient input from\n'
         'the communities it impacts most." — WA State Board of Ed, Dec 2025',
         5.65, 6.56, 7.3, 0.65,
         size=9.5, italic=True, color=SLATE, wrap=True)

# PathFinder logo mark top-right
add_rect(s1, 12.6, 0.2, 0.5, 0.5, fill=GREEN)
add_text(s1, "PF", 12.6, 0.22, 0.5, 0.46,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — LIVE DEMO
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
bg_rect(s2, BG)

add_label(s2, "LIVE DEMO  ·  PATHFINDER AGENT", 0.5, 0.3, color=GREEN)
add_text(s2, "One sentence in.\nA career path out.", 0.5, 0.65, 7, 1.5,
         size=36, bold=True, color=DARK, font="Calibri")

# Demo prompt box
add_rect(s2, 0.5, 2.1, 8.8, 0.85, fill=WHITE, line_color=GREEN)
add_text(s2, "💬  \"I'm 17, Latina, interested in coding, South Seattle\"",
         0.68, 2.2, 8.5, 0.6, size=13, color=DARK, italic=True)

# Arrow
add_text(s2, "↓  PathFinder generates in < 8 seconds", 0.5, 3.05, 5, 0.35,
         size=11, color=GREEN, font="Courier New")

# 4 roadmap steps
steps = [
    ("Week 1–2",  "Apply to Ada Build — Self-Paced",        "Free · Zero prerequisites · Online"),
    ("Week 3–4",  "Complete Ada Build, apply to Ada Build Live", "6-week evening cohort · 30 seats"),
    ("Month 2",   "Apply to Ada Core — Cohort 22",           "Deadline Jul 1 · 60 seats · Full-time"),
    ("Month 3",   "Ada Core starts → paid internship",        "Amazon · Microsoft · Expedia · $95K–$120K"),
]
step_colors = [GREEN, PURPLE, AMBER, RED]
for i, (week, action, detail) in enumerate(steps):
    t = 3.5 + i * 0.83
    add_rect(s2, 0.5, t, 0.32, 0.32, fill=step_colors[i])
    add_text(s2, str(i+1), 0.5, t+0.01, 0.32, 0.3,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, week,   0.95, t,      1.5,  0.2, size=9,  color=step_colors[i], font="Courier New")
    add_text(s2, action, 0.95, t+0.18, 5,    0.28, size=12, bold=True, color=DARK)
    add_text(s2, detail, 0.95, t+0.46, 5,    0.25, size=10, color=MUTED)

# Right panel — mentor + job card
add_rect(s2, 9.6, 0.5, 3.5, 2.8, fill=PURPLE_LIGHT, line_color=PURPLE)
add_text(s2, "MATCHED MENTOR", 9.75, 0.6, 3.2, 0.25,
         size=8, bold=True, color=PURPLE, font="Courier New")
add_text(s2, "Sarah Reyes", 9.75, 0.88, 3.2, 0.38,
         size=16, bold=True, color=DARK, font="Calibri")
add_text(s2, "Software Engineer II · Microsoft\nAda Core Cohort 14 · South Seattle",
         9.75, 1.25, 3.2, 0.55, size=10, color=SLATE, wrap=True)
add_text(s2, '"I was exactly where you are.\nAda changed everything — and\nit\'s genuinely free."',
         9.75, 1.82, 3.2, 0.9, size=10.5, italic=True, color=PURPLE, wrap=True)

add_rect(s2, 9.6, 3.5, 3.5, 2.0, fill=GREEN_LIGHT, line_color=GREEN)
add_text(s2, "WHERE THIS LEADS", 9.75, 3.6, 3.2, 0.25,
         size=8, bold=True, color=GREEN, font="Courier New")
add_text(s2, "Junior Software\nEngineer", 9.75, 3.88, 3.2, 0.65,
         size=18, bold=True, color=DARK, font="Calibri")
add_text(s2, "$95,000 – $120,000", 9.75, 4.52, 3.2, 0.35,
         size=14, bold=True, color=GREEN, font="Calibri")
add_text(s2, "Amazon · Microsoft · Expedia\n12–18 months from today",
         9.75, 4.88, 3.2, 0.5, size=10, color=SLATE, wrap=True)

add_text(s2, "02 / 04", 0.5, 7.1, 1.5, 0.3, size=9,
         color=MUTED, font="Courier New")
add_rect(s2, 12.6, 0.2, 0.5, 0.5, fill=GREEN)
add_text(s2, "PF", 12.6, 0.22, 0.5, 0.46,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW IT WORKS + GOV INTEL
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
bg_rect(s3, BG)

add_label(s3, "TECHNICAL ARCHITECTURE  ·  GOVERNMENT INTELLIGENCE LAYER", 0.5, 0.28, color=PURPLE)
add_text(s3, "Not a chatbot wrapper.\nA provenance-aware multi-agent system.", 0.5, 0.6, 9, 1.0,
         size=26, bold=True, color=DARK, font="Calibri")

# Pipeline nodes
nodes = [
    ("🗣️", "Intake\nAgent",     "Claude Haiku",  True),
    ("🗺️", "Resource\nMatcher", "OpenClaw Tool", False),
    ("🧩", "Gap\nAnalyzer",     "Strands Agent", False),
    ("📋", "Roadmap\nGen",      "Claude Sonnet", False),
    ("🤝", "Mentor\nHandoff",   "HITL Interrupt",False),
]
node_l = 0.5
for i, (icon, name, tech, active) in enumerate(nodes):
    pipe_node(s3, node_l + i * 1.82, 1.75, icon, name, tech, active)
    if i < len(nodes) - 1:
        add_text(s3, "→", node_l + i * 1.82 + 1.47, 2.2, 0.33, 0.35,
                 size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Safety callout
add_rect(s3, 0.5, 3.15, 8.5, 0.65, fill=AMBER_LIGHT, line_color=AMBER)
add_text(s3, "⚠  Safety control: LLM never invents program names, deadlines, or eligibility rules. "
             "All facts come from a curated WA program database. AI synthesizes — data governs.",
         0.65, 3.22, 8.2, 0.5, size=10, color=RGBColor(0x78, 0x35, 0x00), wrap=True)

# Gov intel section
add_rect(s3, 0.5, 4.0, 12.6, 0.01, fill=BORDER)
add_label(s3, "GOVERNMENT INTELLIGENCE LAYER  ·  WHAT YOUTH NEED THAT THE SYSTEM DOESN'T PROVIDE", 0.5, 4.1, color=PURPLE)

gov_items = [
    (PURPLE_LIGHT, PURPLE, "📡 Live Reddit Signal",
     "Real posts from youth seeking help —\ncategorized by need type, urgency, WA location"),
    (RGBColor(0xFF,0xF7,0xF7), RED, "🚨 Gap Detection",
     "Posts where youth need is urgent\nbut no WA program exists to answer it"),
    (GREEN_LIGHT, GREEN, "🤝 HITL Dashboard",
     "AI handled 80% of intake queries.\n1 mentor now serves 10× the caseload"),
    (AMBER_LIGHT, AMBER, "📊 Theme Analysis",
     "Aggregated pain themes for policy:\nchildcare, no-degree, immigration eligibility"),
]

for i, (bg, bc, title, desc) in enumerate(gov_items):
    l = 0.5 + i * 3.18
    add_rect(s3, l, 4.45, 3.0, 2.7, fill=bg, line_color=bc)
    add_text(s3, title, l+0.15, 4.58, 2.7, 0.4, size=11, bold=True, color=bc, wrap=True)
    add_text(s3, desc,  l+0.15, 4.98, 2.7, 0.95, size=10, color=SLATE, wrap=True)

add_text(s3, "03 / 04", 0.5, 7.1, 1.5, 0.3, size=9,
         color=MUTED, font="Courier New")
add_rect(s3, 12.6, 0.2, 0.5, 0.5, fill=GREEN)
add_text(s3, "PF", 12.6, 0.22, 0.5, 0.46,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — IMPACT + CLOSE
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
bg_rect(s4, DARK)

# Left — dark close
add_text(s4, "Washington is rewriting\ngraduation requirements\nright now.", 0.5, 0.5, 6.5, 2.2,
         size=32, bold=True, color=WHITE, font="Calibri")
add_text(s4, "We built the navigation\nlayer they forgot to include.", 0.5, 2.7, 6.5, 1.4,
         size=32, bold=True, color=GREEN, font="Calibri")

add_text(s4,
         "Ada has graduated 1,200 engineers since 2013.\n"
         "PathFinder finds the next 12,000 —\n"
         "and shows policymakers exactly where to fund\n"
         "the programs that don't exist yet.",
         0.5, 4.3, 6.2, 1.8,
         size=14, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

# WA alignment strip
add_rect(s4, 0.5, 6.3, 6.2, 0.9, fill=RGBColor(0x1E, 0x29, 0x3B))
align_items = ["FutureReady 2026", "WSAC BIPOC Gap", "Ada AI Curriculum", "School's Out WA", "UW GEAR UP"]
align_colors = [GREEN, PURPLE, RED, AMBER, RGBColor(0x08, 0x91, 0xB2)]
for i, (item, col) in enumerate(zip(align_items, align_colors)):
    x = 0.65 + i * 1.22
    add_rect(s4, x, 6.42, 0.08, 0.28, fill=col)
    add_text(s4, item, x+0.13, 6.41, 1.05, 0.35, size=8, color=RGBColor(0xCB,0xD5,0xE1), wrap=True)

# Right — rubric proof tiles
add_label(s4, "RUBRIC PROOF", 7.1, 0.35, color=PURPLE)

rubric = [
    (GREEN,  "Human Impact",
     "Marisol: confused → roadmap → mentor → $95K career in < 10 min.\n1 mentor now serves 10× more youth."),
    (PURPLE, "Security / Responsibility",
     "No PII stored. No account required. HITL interrupt for undocumented,\nfoster youth, and emotional distress cases. AI knows its limits."),
    (AMBER,  "Effectiveness",
     "25 verified WA programs. < 8 second roadmap on Groq free tier.\n7 policy gap signals detected from live Reddit data. Tested today."),
    (RED,    "Feasibility",
     "FastAPI + LangGraph + static HTML. Zero infrastructure cost.\nHand to Ada's program team Monday — live same day."),
]

for i, (col, title, body) in enumerate(rubric):
    t = 0.65 + i * 1.62
    add_rect(s4, 7.1, t, 5.9, 1.5, fill=RGBColor(0x1E, 0x29, 0x3B),
             line_color=col)
    add_rect(s4, 7.1, t, 0.18, 1.5, fill=col)
    add_text(s4, title, 7.38, t+0.15, 5.4, 0.35,
             size=13, bold=True, color=col, font="Calibri")
    add_text(s4, body,  7.38, t+0.52, 5.4, 0.9,
             size=10, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

add_text(s4, "04 / 04", 0.5, 7.1, 1.5, 0.3, size=9,
         color=MUTED, font="Courier New")
add_rect(s4, 12.6, 0.2, 0.5, 0.5, fill=GREEN)
add_text(s4, "PF", 12.6, 0.22, 0.5, 0.46,
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────
out = "PathFinder_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
